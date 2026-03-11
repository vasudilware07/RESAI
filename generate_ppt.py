"""
generate_ppt.py - PowerPoint Presentation Generator
======================================================
Generates a 20-slide presentation covering:
  - Project overview & architecture
  - Dataset description
  - Model training & results
  - Fairness & bias detection
  - Explainability techniques
  - Privacy preservation
  - Ethics audit
  - Conclusion & references

Run: python generate_ppt.py  (after running main.py)
"""

import os
import json
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
PLOTS_DIR = os.path.join(BASE_DIR, "outputs", "plots")
REPORTS_DIR = os.path.join(BASE_DIR, "outputs", "reports")
OUTPUT_PPT = os.path.join(BASE_DIR, "outputs", "Responsible_AI_Hiring_Presentation.pptx")


# ─── Color palette ───────────────────────────────────────────────────────────
BG_DARK = RGBColor(0x1B, 0x1B, 0x2F)   # Dark navy
BG_LIGHT = RGBColor(0xF5, 0xF5, 0xF5)  # Light gray
ACCENT = RGBColor(0x00, 0x96, 0xC7)     # Blue accent
ACCENT2 = RGBColor(0xE7, 0x4C, 0x3C)    # Red accent
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x20, 0x20, 0x20)
GRAY = RGBColor(0x60, 0x60, 0x60)


def load_json_file(filename):
    path = os.path.join(REPORTS_DIR, filename)
    if os.path.exists(path):
        with open(path) as f:
            return json.load(f)
    return None


def add_background(slide, color=BG_DARK):
    """Set slide background color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_slide(prs, title, subtitle=""):
    """Add a styled title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    add_background(slide, BG_DARK)

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(8.4), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    if subtitle:
        txBox2 = slide.shapes.add_textbox(Inches(1.0), Inches(3.6), Inches(8.0), Inches(1.0))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(18)
        p2.font.color.rgb = ACCENT
        p2.alignment = PP_ALIGN.CENTER

    return slide


def add_content_slide(prs, title, bullets, image_path=None):
    """Add a slide with title, bullet points, and optional image."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    add_background(slide, BG_LIGHT)

    # Title bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.9))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_DARK
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Inches(0.5)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Content area
    if image_path and os.path.exists(os.path.join(PLOTS_DIR, image_path)):
        # Left bullets, right image
        txt_left = Inches(0.5)
        txt_width = Inches(4.2)
        img_left = Inches(4.8)
        img_top = Inches(1.2)
        img_width = Inches(5.0)

        txBox = slide.shapes.add_textbox(txt_left, Inches(1.2), txt_width, Inches(5.5))
        tf2 = txBox.text_frame
        tf2.word_wrap = True
        for bullet in bullets:
            p = tf2.add_paragraph()
            p.text = bullet
            p.font.size = Pt(13)
            p.font.color.rgb = BLACK
            p.space_after = Pt(6)
            p.level = 0

        # Remove first empty paragraph
        if tf2.paragraphs[0].text == "":
            tf2.paragraphs[0].text = bullets[0] if bullets else ""

        try:
            slide.shapes.add_picture(
                os.path.join(PLOTS_DIR, image_path),
                img_left, img_top, img_width
            )
        except Exception:
            pass
    else:
        # Full-width bullets
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9.0), Inches(5.5))
        tf2 = txBox.text_frame
        tf2.word_wrap = True
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf2.paragraphs[0]
            else:
                p = tf2.add_paragraph()
            p.text = bullet
            p.font.size = Pt(14)
            p.font.color.rgb = BLACK
            p.space_after = Pt(8)

    return slide


def add_image_slide(prs, title, image_path, caption=""):
    """Slide with full-width image."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, BG_LIGHT)

    # Title bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.9))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_DARK
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Inches(0.5)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE

    full_path = os.path.join(PLOTS_DIR, image_path)
    if os.path.exists(full_path):
        slide.shapes.add_picture(full_path, Inches(0.5), Inches(1.1), Inches(9.0))
    else:
        txBox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1))
        txBox.text_frame.paragraphs[0].text = f"[Image not found: {image_path}. Run main.py first.]"

    if caption:
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.5))
        p = txBox.text_frame.paragraphs[0]
        p.text = caption
        p.font.size = Pt(11)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER

    return slide


def generate_presentation():
    """Generate the full 20-slide presentation."""
    print("Generating PowerPoint Presentation...")
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ── SLIDE 1: Title ───────────────────────────────────────────────
    add_title_slide(
        prs,
        "Bias Detection & Responsible AI\nin AI Hiring Systems",
        "Academic Project | IBM HR Analytics Dataset"
    )

    # ── SLIDE 2: Agenda ─────────────────────────────────────────────
    add_content_slide(prs, "Agenda", [
        "1. Introduction to Responsible AI",
        "2. Problem Statement & Dataset",
        "3. System Architecture",
        "4. Data Preprocessing & Feature Engineering",
        "5. AI Hiring Model – Training & Evaluation",
        "6. Fairness & Bias Detection (Unit 2)",
        "7. Bias Visualization & Analysis",
        "8. Interpretability & Explainability (Unit 3)",
        "9. SHAP Analysis & Feature Importance",
        "10. Individual Candidate Explanations",
        "11. Ethics & Accountability Audit (Unit 4)",
        "12. Privacy Preservation (Unit 5)",
        "13. Privacy-Accuracy-Fairness Tradeoffs",
        "14. Correlation vs Causation Analysis",
        "15. Nudge Theory & Counterfactual Explanations",
        "16. Dashboard Demo",
        "17. Key Findings & Conclusion",
    ])

    # ── SLIDE 3: Introduction to RAI ─────────────────────────────────
    add_content_slide(prs, "Introduction to Responsible AI (Unit 1)", [
        "• Responsible AI ensures AI systems are fair, transparent,",
        "  accountable, and privacy-preserving.",
        "",
        "• Key Principles:",
        "  → Fairness: No discrimination based on protected attributes",
        "  → Transparency: Explainable decision-making",
        "  → Accountability: Clear ownership of AI outcomes",
        "  → Privacy: Protection of personal data",
        "  → Beneficence: AI should benefit all stakeholders",
        "",
        "• Why it matters in Hiring:",
        "  → AI hiring tools can perpetuate historical biases",
        "  → Legal compliance (EU AI Act, EEOC 4/5ths rule)",
        "  → Ethical obligation to candidates",
    ])

    # ── SLIDE 4: Problem Statement ───────────────────────────────────
    add_content_slide(prs, "Problem Statement", [
        "Problem:",
        "AI-powered hiring systems can introduce or amplify biases",
        "against protected groups (gender, age, ethnicity).",
        "",
        "Objective:",
        "Build a Responsible AI pipeline that:",
        "  ✓ Detects bias in hiring predictions",
        "  ✓ Explains model decisions transparently",
        "  ✓ Audits ethical risks systematically",
        "  ✓ Preserves candidate privacy",
        "",
        "Dataset: IBM HR Employee Attrition",
        "  • 1,470 employees × 35 features",
        "  • Target: Attrition (Yes/No → Hired/Not Hired)",
        "  • Sensitive attributes: Gender, Age",
    ])

    # ── SLIDE 5: Architecture ────────────────────────────────────────
    add_content_slide(prs, "System Architecture", [
        "Pipeline Flow:",
        "",
        "  Dataset (IBM HR)        ",
        "       ↓                  ",
        "  Data Preprocessing      ",
        "       ↓                  ",
        "  Model Training (LR, RF, DT)",
        "       ↓                  ",
        "  ┌──────┬──────┬──────┐ ",
        "  ↓      ↓      ↓      ↓ ",
        "  Fairness  Explain  Privacy",
        "  Analysis  ability  Module ",
        "  ↓      ↓      ↓        ",
        "  └──────┼──────┘        ",
        "         ↓               ",
        "  Ethics Audit → Dashboard",
        "",
        "Tech Stack:",
        "  Python, Scikit-learn, SHAP, Fairlearn,",
        "  diffprivlib, Streamlit, python-pptx",
    ])

    # ── SLIDE 6: Data Preprocessing ──────────────────────────────────
    add_content_slide(prs, "Data Preprocessing & Feature Engineering", [
        "Steps performed:",
        "",
        "1. Data Cleaning:",
        "   • Dropped constant columns (EmployeeCount, Over18, etc.)",
        "   • Mapped target: Attrition Yes→1, No→0",
        "   • Handled missing values",
        "",
        "2. Feature Engineering:",
        "   • Created AgeGroup bins (18-30, 31-40, 41-50, 51-60)",
        "   • Created IncomeGroup quartiles",
        "   • Binary gender encoding",
        "",
        "3. Encoding & Scaling:",
        "   • Label-encoded all categorical features",
        "   • StandardScaler on numeric features",
        "",
        "4. Sensitive Attributes Identified:",
        "   • Gender (Male/Female)",
        "   • Age Group (Young/Senior)",
    ])

    # ── SLIDE 7: Model Training ──────────────────────────────────────
    model_metrics = load_json_file("model_metrics.json")
    bullets = [
        "Three models trained for hiring prediction:",
        "",
        "1. Logistic Regression (interpretable, baseline)",
        "2. Random Forest (ensemble, high accuracy)",
        "3. Decision Tree (intrinsically interpretable)",
        "",
        "Evaluation Metrics:",
    ]
    if model_metrics:
        for name, metrics in model_metrics.items():
            m = metrics if isinstance(metrics, dict) else {}
            bullets.append(
                f"  {name}: Acc={m.get('Accuracy','?')}, "
                f"P={m.get('Precision','?')}, R={m.get('Recall','?')}, "
                f"F1={m.get('F1 Score','?')}"
            )

    add_content_slide(prs, "AI Hiring Model – Training & Evaluation", bullets,
                      "model_comparison.png")

    # ── SLIDE 8: Confusion Matrices ──────────────────────────────────
    add_image_slide(prs, "Model Evaluation – Confusion Matrices",
                    "confusion_matrices.png",
                    "Confusion matrices showing True/False Positive/Negative rates for each model")

    # ── SLIDE 9: Fairness Metrics ────────────────────────────────────
    add_content_slide(prs, "Fairness & Bias Detection (Unit 2)", [
        "Fairness Metrics Implemented:",
        "",
        "1. Demographic Parity:",
        "   P(ŷ=1|G=0) should equal P(ŷ=1|G=1)",
        "",
        "2. Equal Opportunity:",
        "   TPR should be equal across groups",
        "",
        "3. Disparate Impact (80% Rule):",
        "   Ratio = P(ŷ=1|unprivileged) / P(ŷ=1|privileged)",
        "   Ratio < 0.8 → adverse impact detected",
        "",
        "4. Statistical Parity Difference:",
        "   SPD = P(ŷ=1|unprivileged) - P(ŷ=1|privileged)",
        "   Ideal value = 0",
        "",
        "Libraries: Fairlearn, AIF360, custom implementations",
    ], "fairness_comparison.png")

    # ── SLIDE 10: Bias Heatmap ───────────────────────────────────────
    add_image_slide(prs, "Bias Heatmap – All Models & Attributes",
                    "bias_heatmap.png",
                    "Heatmap showing fairness metrics across models for Gender and Age groups")

    # ── SLIDE 11: Group Fairness ─────────────────────────────────────
    add_image_slide(prs, "Group Fairness – Gender Analysis",
                    "group_fairness_gender.png",
                    "Positive prediction rates by gender for each model")

    # ── SLIDE 12: Explainability Overview ────────────────────────────
    add_content_slide(prs, "Interpretability & Explainability (Unit 3)", [
        "Two categories of explainability:",
        "",
        "A. Intrinsic Methods (built-in interpretability):",
        "   • Decision Tree – directly readable rules",
        "   • Logistic Regression – feature coefficients",
        "",
        "B. Post-hoc Methods (explain after training):",
        "   • SHAP (SHapley Additive exPlanations)",
        "     → Game-theoretic feature attribution",
        "   • Feature Importance (Random Forest)",
        "     → Gini / permutation importance",
        "   • Partial Dependence Plots",
        "     → Marginal effect of each feature",
        "",
        "Why it matters:",
        "   • GDPR Article 22: right to explanation",
        "   • Candidates deserve to know WHY",
    ], "feature_importance.png")

    # ── SLIDE 13: SHAP Analysis ──────────────────────────────────────
    add_image_slide(prs, "SHAP Summary Plot – Feature Impact",
                    "shap_summary.png",
                    "Each point = one employee. Color = feature value. X-axis = impact on prediction")

    # ── SLIDE 14: Individual Explanation ─────────────────────────────
    add_content_slide(prs, "Individual Candidate Explanation", [
        "For each candidate, SHAP explains:",
        "  → Which features pushed toward 'Hired'",
        "  → Which features pushed toward 'Not Hired'",
        "",
        "Example: Candidate #0",
        "  • OverTime → strong push toward Not Hired",
        "  • MonthlyIncome → push toward Hired",
        "  • TotalWorkingYears → push toward Hired",
        "",
        "This provides:",
        "  ✓ Transparency for the candidate",
        "  ✓ Actionable feedback",
        "  ✓ Audit trail for the organization",
    ], "individual_explanation.png")

    # ── SLIDE 15: Decision Tree ──────────────────────────────────────
    add_image_slide(prs, "Decision Tree Visualization (Intrinsic)",
                    "decision_tree.png",
                    "The Decision Tree model is inherently interpretable – each path is a hiring rule")

    # ── SLIDE 16: Ethics Audit ───────────────────────────────────────
    risks = load_json_file("ethical_risks.json")
    bullets = [
        "Automated AI Audit Report includes:",
        "",
        "1. Bias findings across all models",
        "2. Ethical risk severity (HIGH / MEDIUM / LOW)",
        "3. Model transparency assessment",
        "4. Regulatory compliance notes",
        "5. Actionable recommendations",
        "",
        "Ethical Principles evaluated:",
        "  • Fairness    • Transparency",
        "  • Accountability  • Privacy",
        "  • Beneficence    • Non-maleficence",
    ]
    if risks:
        high = sum(1 for r in risks if r.get("Risk") == "HIGH")
        med = sum(1 for r in risks if r.get("Risk") == "MEDIUM")
        low = sum(1 for r in risks if r.get("Risk") == "LOW")
        bullets.append(f"")
        bullets.append(f"Risk Summary: {high} HIGH, {med} MEDIUM, {low} LOW")

    add_content_slide(prs, "Ethics & Accountability Audit (Unit 4)", bullets,
                      "risk_summary.png")

    # ── SLIDE 17: Privacy ────────────────────────────────────────────
    privacy = load_json_file("privacy_results.json")
    bullets = [
        "Technique: Differential Privacy (ε-DP)",
        "",
        "• Adds calibrated noise to model training",
        "• Privacy budget ε controls the tradeoff:",
        "  → Low ε = strong privacy, lower accuracy",
        "  → High ε = weak privacy, higher accuracy",
        "",
        "Implementation: IBM diffprivlib",
        "  • DP Logistic Regression at ε = 0.01 to 50",
        "",
    ]
    if privacy and "baseline" in privacy:
        b = privacy["baseline"]
        bullets.append(f"Baseline: Acc={b['accuracy']}, F1={b['f1_score']}")
        if privacy.get("comparison"):
            best = max(privacy["comparison"], key=lambda x: x["accuracy"])
            worst = min(privacy["comparison"], key=lambda x: x["accuracy"])
            bullets.append(f"Best DP:  ε={best['epsilon']}, Acc={best['accuracy']}")
            bullets.append(f"Most Private: ε={worst['epsilon']}, Acc={worst['accuracy']}")

    add_content_slide(prs, "Privacy Preservation (Unit 5)", bullets,
                      "privacy_accuracy_tradeoff.png")

    # ── SLIDE 18: Privacy Tradeoffs ──────────────────────────────────
    add_image_slide(prs, "Accuracy vs Fairness at Different Privacy Levels",
                    "accuracy_vs_fairness_privacy.png",
                    "Each point = a DP model at different ε. Star = baseline without privacy")

    # ── SLIDE 19: Dashboard Screenshots ──────────────────────────────
    add_content_slide(prs, "Responsible AI Dashboard", [
        "Interactive Streamlit Dashboard features:",
        "",
        "📊 Dataset Overview:",
        "  • Employee statistics, distributions, demographics",
        "",
        "🤖 Model Performance:",
        "  • Side-by-side metric comparison",
        "  • Confusion matrices",
        "",
        "⚖️ Fairness & Bias:",
        "  • Interactive fairness metric tables",
        "  • Bias heatmaps and comparison charts",
        "",
        "🔍 Explainability:",
        "  • SHAP plots, feature importance, PDP",
        "",
        "📋 Ethics Audit:",
        "  • Full audit report with risk table",
        "",
        "🔗 Correlation & Causation:",
        "  • Correlation matrix, proxy variables, Simpson's Paradox",
        "",
        "🎯 Nudge Theory:",
        "  • Threshold nudging, counterfactual explanations",
        "",
        "🔒 Privacy Analysis:",
        "  • DP comparison tables and tradeoff plots",
        "",
        "Run: streamlit run dashboard.py",
    ])

    # ── SLIDE 20: Correlation vs Causation ────────────────────────────
    ccn_results = load_json_file("correlation_causation_nudge.json")
    corr_bullets = [
        "Correlation ≠ Causation in AI Hiring:",
        "",
        "Correlation Analysis:",
        "  • Full feature correlation matrix computed",
        "  • Identified top features correlated with Attrition",
        "  • VIF analysis to detect multicollinearity",
        "",
        "Causation Analysis:",
        "  • Simpson's Paradox detection across subgroups",
        "  • Proxy variable identification (features correlated",
        "    with protected attributes like Gender, Age)",
        "  • Confounding variable analysis using partial correlations",
        "",
        "Key Finding:",
        "  Many features that correlate with attrition also correlate",
        "  with protected attributes → potential indirect discrimination",
    ]
    if ccn_results and "causation" in ccn_results:
        caus = ccn_results["causation"]
        if "proxy_variables" in caus and caus["proxy_variables"]:
            proxies = caus["proxy_variables"][:5]
            names = [p["proxy_feature"] if isinstance(p, dict) else str(p) for p in proxies]
            corr_bullets.append(f"  Proxy variables found: {', '.join(names)}")
    add_content_slide(prs, "Correlation vs Causation Analysis", corr_bullets,
                      "target_correlation_ranking.png")

    # ── SLIDE 21: Simpson's Paradox & Proxies ────────────────────────
    add_image_slide(prs, "Simpson's Paradox & Proxy Variable Detection",
                    "proxy_variables.png",
                    "Proxy variables can transmit bias from protected attributes to model predictions")

    # ── SLIDE 22: Nudge Theory ───────────────────────────────────────
    nudge_bullets = [
        "Nudge Theory (Thaler & Sunstein, 2008):",
        "  Small design changes guide better decisions",
        "  without restricting freedom of choice.",
        "",
        "Applied to AI Hiring:",
        "",
        "1. Threshold Nudging:",
        "   • Default threshold = 0.5 may be unfair",
        "   • Optimized threshold balances accuracy & fairness",
        "   • Small shift → significant fairness improvement",
        "",
        "2. Counterfactual Explanations:",
        "   • 'What minimal change would flip the decision?'",
        "   • Provides actionable feedback to candidates",
        "   • Ensures transparency & right to explanation",
    ]
    if ccn_results and "nudge" in ccn_results:
        nudge = ccn_results["nudge"]
        if "optimal_threshold" in nudge:
            opt = nudge["optimal_threshold"]
            nudge_bullets.append(f"")
            nudge_bullets.append(f"  Optimal threshold: {opt.get('threshold', 'N/A')}")
    add_content_slide(prs, "Nudge Theory & Counterfactual Explanations", nudge_bullets,
                      "threshold_nudge.png")

    # ── SLIDE 23: Conclusion ─────────────────────────────────────────
    add_title_slide(prs, "Conclusion & Key Takeaways", "")

    # Add conclusion content on the same dark slide
    slide = prs.slides[-1]
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(3.8), Inches(8.4), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    conclusions = [
        "✓ Built a complete Responsible AI pipeline for hiring",
        "✓ Detected measurable bias across Gender and Age groups",
        "✓ Provided model explanations using SHAP & Decision Trees",
        "✓ Generated automated ethics audit with risk assessment",
        "✓ Demonstrated privacy–accuracy–fairness tradeoffs",
        "✓ Analyzed correlation vs causation & proxy variables",
        "✓ Applied nudge theory & counterfactual explanations",
        "✓ Created interactive dashboard for stakeholder review",
        "",
        "Key Insight: Responsible AI is not optional –",
        "it is essential for trustworthy AI deployment.",
    ]
    for i, text in enumerate(conclusions):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.space_after = Pt(4)

    # ── Save ─────────────────────────────────────────────────────────
    os.makedirs(os.path.dirname(OUTPUT_PPT), exist_ok=True)
    prs.save(OUTPUT_PPT)
    print(f"\n✓ Presentation saved: {OUTPUT_PPT}")
    print(f"  Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    generate_presentation()
